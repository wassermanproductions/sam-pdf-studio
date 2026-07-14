#!/usr/bin/env python3
from __future__ import annotations

import argparse
import contextlib
import io
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
from pathlib import Path
from typing import Any

ENGINE_VERSION = "0.2.0"


def respond(ok: bool, **payload: Any) -> int:
    print(json.dumps({"ok": ok, **payload}, indent=2, sort_keys=True))
    return 0 if ok else 1


def ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def safe_stem(path: str | Path) -> str:
    stem = Path(path).stem
    cleaned = "".join(ch if ch.isalnum() or ch in "-_" else "-" for ch in stem)
    while "--" in cleaned:
        cleaned = cleaned.replace("--", "-")
    return cleaned.strip("-") or "output"


def parse_page_spec(spec: str | None, total_pages: int) -> list[int]:
    if spec is None or spec.strip().lower() in {"", "all"}:
        return list(range(total_pages))

    pages: set[int] = set()
    for chunk in spec.split(","):
        part = chunk.strip()
        if not part:
            continue
        if "-" in part:
            start_text, end_text = part.split("-", 1)
            start = int(start_text)
            end = int(end_text)
            if start > end:
                start, end = end, start
            for page in range(start, end + 1):
                pages.add(page - 1)
        else:
            pages.add(int(part) - 1)

    invalid = [page + 1 for page in pages if page < 0 or page >= total_pages]
    if invalid:
        raise ValueError(f"page selection out of range: {invalid}; document has {total_pages} pages")
    return sorted(pages)


def search_terms(text: str) -> list[str]:
    terms: list[str] = []
    raw = text.strip()
    collapsed = " ".join(raw.split())
    for term in [raw, collapsed, *[line.strip() for line in raw.splitlines()]]:
        if term and term not in terms:
            terms.append(term)
    return terms


def search_rects(page: Any, text: str) -> list[Any]:
    for term in search_terms(text):
        rects = page.search_for(term)
        if rects:
            return rects
    return []


def parse_rect(rect_text: str | None) -> tuple[float, float, float, float] | None:
    if not rect_text:
        return None
    parts = [float(part.strip()) for part in rect_text.split(",")]
    if len(parts) != 4:
        raise ValueError("--rect must be x0,y0,x1,y1")
    return (min(parts[0], parts[2]), min(parts[1], parts[3]), max(parts[0], parts[2]), max(parts[1], parts[3]))


def mirrored_vertical_rect(page: Any, rect: Any) -> Any:
    import fitz

    return fitz.Rect(
        rect.x0,
        page.rect.y1 - (rect.y1 - page.rect.y0),
        rect.x1,
        page.rect.y1 - (rect.y0 - page.rect.y0),
    )


def rects_intersect(a: Any, b: Any) -> bool:
    return not (a.x1 < b.x0 or a.x0 > b.x1 or a.y1 < b.y0 or a.y0 > b.y1)


def rect_overlap_area(a: Any, b: Any) -> float:
    x_overlap = max(0, min(a.x1, b.x1) - max(a.x0, b.x0))
    y_overlap = max(0, min(a.y1, b.y1) - max(a.y0, b.y0))
    return x_overlap * y_overlap


def rect_proximity(anchor: Any, candidate: Any) -> float:
    overlap = rect_overlap_area(anchor, candidate)
    if overlap > 0:
        return -overlap
    dx = max(anchor.x0 - candidate.x1, candidate.x0 - anchor.x1, 0)
    dy = max(anchor.y0 - candidate.y1, candidate.y0 - anchor.y1, 0)
    return dx * dx + dy * dy


def nearest_search_rect(page: Any, click_rect: Any, text: str) -> Any | None:
    matches = search_rects(page, text)
    if not matches:
        return None
    anchors = [click_rect, mirrored_vertical_rect(page, click_rect)]
    return min(matches, key=lambda match: min(rect_proximity(anchor, match) for anchor in anchors))


def pdf_color_to_rgb(value: int | None) -> tuple[float, float, float]:
    if value is None:
        return (0, 0, 0)
    return (((value >> 16) & 255) / 255.0, ((value >> 8) & 255) / 255.0, (value & 255) / 255.0)


def builtin_font_alias(font_name: str | None) -> str:
    name = (font_name or "").lower()
    if "courier" in name or "mono" in name:
        if "bold" in name and ("italic" in name or "oblique" in name):
            return "cobo"
        if "bold" in name:
            return "cob"
        if "italic" in name or "oblique" in name:
            return "coi"
        return "cour"
    if "times" in name or "serif" in name:
        if "bold" in name and ("italic" in name or "oblique" in name):
            return "tibi"
        if "bold" in name:
            return "tibo"
        if "italic" in name or "oblique" in name:
            return "tiit"
        return "tiro"
    if "bold" in name and ("italic" in name or "oblique" in name):
        return "hebo"
    if "bold" in name:
        return "hebo"
    if "italic" in name or "oblique" in name:
        return "heit"
    return "helv"


def align_constant(name: str | None) -> int:
    """Map a friendly alignment name to a fitz TEXT_ALIGN_* constant."""
    import fitz

    return {
        "center": fitz.TEXT_ALIGN_CENTER,
        "right": fitz.TEXT_ALIGN_RIGHT,
    }.get((name or "left").strip().lower(), fitz.TEXT_ALIGN_LEFT)


def draw_text_underlines(page: Any, clip: Any, color: Any, width: float = 0.6) -> None:
    """Underline every text line whose center sits inside `clip`, drawing a
    thin rule just below each line's baseline. Additive vector drawing only —
    it never touches or removes the glyphs, so text is preserved."""
    import fitz

    clip_rect = fitz.Rect(clip)
    stroke = color or (0, 0, 0)
    for block in page.get_text("dict", clip=clip_rect).get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = [s for s in line.get("spans", []) if (s.get("text") or "").strip()]
            if not spans:
                continue
            bbox = line.get("bbox")
            if not bbox:
                continue
            center = fitz.Point((bbox[0] + bbox[2]) / 2, (bbox[1] + bbox[3]) / 2)
            if not clip_rect.contains(center):
                continue
            y = bbox[3] - 1
            page.draw_line(fitz.Point(bbox[0], y), fitz.Point(bbox[2], y), color=stroke, width=width)


def style_for_rect(page: Any, rect: Any) -> dict[str, Any]:
    data = page.get_text("dict")
    best_span: dict[str, Any] | None = None
    best_area = 0.0
    for block in data.get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                bbox = span.get("bbox")
                if not bbox:
                    continue
                try:
                    import fitz

                    span_rect = fitz.Rect(bbox)
                except Exception:
                    continue
                if not rects_intersect(span_rect, rect):
                    continue
                x_overlap = max(0, min(span_rect.x1, rect.x1) - max(span_rect.x0, rect.x0))
                y_overlap = max(0, min(span_rect.y1, rect.y1) - max(span_rect.y0, rect.y0))
                area = x_overlap * y_overlap
                if area > best_area:
                    best_span = span
                    best_area = area

    if not best_span:
        return {
            "fontname": "helv",
            "fontsize": max(8, rect.height * 0.72),
            "color": (0, 0, 0),
            "rawfont": None,
            "origin": None,
        }

    return {
        "fontname": builtin_font_alias(best_span.get("font")),
        "fontsize": float(best_span.get("size") or max(8, rect.height * 0.72)),
        "color": pdf_color_to_rgb(best_span.get("color")),
        "rawfont": best_span.get("font"),
        "origin": tuple(best_span.get("origin") or ()) or None,
    }


def strip_subset_prefix(font_name: str) -> str:
    # Embedded subset fonts are named like "ABCDEF+Arial-BoldMT".
    if len(font_name) > 7 and font_name[6] == "+" and font_name[:6].isalpha() and font_name[:6].isupper():
        return font_name[7:]
    return font_name


def font_covers_text(fontbuffer: bytes, text: str) -> bool:
    import fitz

    try:
        font = fitz.Font(fontbuffer=fontbuffer)
    except Exception:
        return False
    # Spaces count too: subset fonts often omit the space glyph, which would
    # render as tofu boxes between words.
    return all(font.has_glyph(ord(ch)) for ch in text if ch not in "\n\r\t")


def font_file_covers_text(path: str, text: str) -> bool:
    import fitz

    try:
        font = fitz.Font(fontfile=path)
    except Exception:
        return False
    return all(font.has_glyph(ord(ch)) for ch in text if ch not in "\n\r\t")


_SYSTEM_FONT_DIRS = (
    "/System/Library/Fonts/Supplemental",
    "/Library/Fonts",
    os.path.expanduser("~/Library/Fonts"),
)


def system_font_file(raw_name: str, text: str) -> str | None:
    """Best-effort lookup of an installed .ttf/.otf matching the PDF font name."""
    cleaned = strip_subset_prefix(raw_name or "")
    if not cleaned:
        return None
    lowered = cleaned.lower()
    is_bold = "bold" in lowered
    is_italic = "italic" in lowered or "oblique" in lowered
    family = re.split(r"[-,]", cleaned)[0]
    family = re.sub(r"(MT|PS|Std|Pro)$", "", family)
    # Split CamelCase into words: "TimesNewRoman" -> "Times New Roman"
    family_words = re.sub(r"(?<=[a-z])(?=[A-Z])", " ", family).strip()
    if not family_words:
        return None

    candidates = []
    base = family_words
    if is_bold and is_italic:
        candidates += [f"{base} Bold Italic", f"{base}BdIt", f"{base}-BoldItalic"]
    elif is_bold:
        candidates += [f"{base} Bold", f"{base}-Bold"]
    elif is_italic:
        candidates += [f"{base} Italic", f"{base}-Italic"]
    candidates.append(base)

    for directory in _SYSTEM_FONT_DIRS:
        if not os.path.isdir(directory):
            continue
        try:
            entries = os.listdir(directory)
        except OSError:
            continue
        by_lower = {entry.lower(): entry for entry in entries}
        for candidate in candidates:
            for ext in (".ttf", ".otf"):
                key = f"{candidate.lower()}{ext}"
                if key in by_lower:
                    path = os.path.join(directory, by_lower[key])
                    if font_file_covers_text(path, text):
                        return path
    return None


def resolve_replacement_font(doc: Any, page: Any, style: dict[str, Any], text: str, scratch: list[str]) -> dict[str, Any]:
    """Pick the closest font for replacement text, matching the original.

    Preference order: the document's own embedded font (when it covers the
    replacement's glyphs) -> an installed system font with the same family ->
    the base-14 alias the old code used.
    Returns kwargs for insert_text/insert_textbox.
    """
    raw_name = style.get("rawfont") or ""
    if raw_name:
        target = raw_name.lower()
        try:
            fonts = page.get_fonts(full=False)
        except Exception:
            fonts = []
        for entry in fonts:
            xref, _ext, _ftype, basefont = entry[0], entry[1], entry[2], entry[3]
            if not basefont:
                continue
            if strip_subset_prefix(basefont).lower() != strip_subset_prefix(raw_name).lower() \
               and target not in basefont.lower():
                continue
            try:
                _name, ext, _t, buffer = doc.extract_font(xref)
            except Exception:
                continue
            if not buffer or ext not in ("ttf", "otf", "cff"):
                continue
            if not font_covers_text(buffer, text):
                continue
            handle, path = tempfile.mkstemp(suffix=f".{ext}")
            with os.fdopen(handle, "wb") as fh:
                fh.write(buffer)
            scratch.append(path)
            return {"fontname": "SamPDFEmbedded", "fontfile": path}

        system_path = system_font_file(raw_name, text)
        if system_path:
            return {"fontname": "SamPDFSystem", "fontfile": system_path}

    return {"fontname": style["fontname"], "fontfile": None}


def replacement_text_length(text: str, font_kwargs: dict[str, Any], fontsize: float) -> float:
    import fitz

    try:
        if font_kwargs.get("fontfile"):
            font = fitz.Font(fontfile=font_kwargs["fontfile"])
        else:
            font = fitz.Font(font_kwargs["fontname"])
        return font.text_length(text, fontsize=fontsize)
    except Exception:
        return len(text) * fontsize * 0.55


@contextlib.contextmanager
def silence_process_output():
    devnull = os.open(os.devnull, os.O_WRONLY)
    old_stdout = os.dup(1)
    old_stderr = os.dup(2)
    try:
        os.dup2(devnull, 1)
        os.dup2(devnull, 2)
        yield
    finally:
        os.dup2(old_stdout, 1)
        os.dup2(old_stderr, 2)
        os.close(old_stdout)
        os.close(old_stderr)
        os.close(devnull)


def import_status(name: str, import_name: str | None = None) -> dict[str, Any]:
    import_name = import_name or name
    try:
        capture = io.StringIO()
        with silence_process_output(), contextlib.redirect_stdout(capture), contextlib.redirect_stderr(capture):
            module = __import__(import_name)
        version = getattr(module, "__version__", None)
        if import_name == "fitz":
            version = getattr(module, "version", ["unknown"])[0]
        return {"ok": True, "version": str(version or "unknown"), "path": getattr(module, "__file__", None), "error": None}
    except Exception as exc:
        return {"ok": False, "version": None, "path": None, "error": str(exc)}


def find_binary(name: str) -> str | None:
    found = shutil.which(name)
    if found:
        return found
    local = Path(sys.executable).resolve().parent / name
    if local.exists():
        return str(local)
    return None


def binary_status(name: str, version_args: list[str] | None = None) -> dict[str, Any]:
    path = find_binary(name)
    if not path:
        return {"ok": False, "version": None, "path": None, "error": "not found"}
    version = "unknown"
    if version_args is not None:
        try:
            result = subprocess.run([path] + version_args, text=True, capture_output=True, timeout=10)
            version_text = (result.stdout or result.stderr).strip().splitlines()
            if version_text:
                version = version_text[0]
        except Exception as exc:
            return {"ok": True, "version": version, "path": path, "error": str(exc)}
    return {"ok": True, "version": version, "path": path, "error": None}


def cmd_health(_args: argparse.Namespace) -> int:
    packages = {
        "pymupdf": import_status("pymupdf", "fitz"),
        "pypdf": import_status("pypdf"),
        "pdfplumber": import_status("pdfplumber"),
        "reportlab": import_status("reportlab"),
        "pillow": import_status("pillow", "PIL"),
        "pdf2docx": import_status("pdf2docx"),
        "ocrmypdf": import_status("ocrmypdf"),
        "pymupdf4llm": import_status("pymupdf4llm"),
        "openpyxl": import_status("openpyxl"),
        "python-pptx": import_status("python-pptx", "pptx"),
    }
    binaries = {
        "qpdf": binary_status("qpdf", ["--version"]),
        "tesseract": binary_status("tesseract", ["--version"]),
        "gs": binary_status("gs", ["--version"]),
        "pdftoppm": binary_status("pdftoppm", ["-v"]),
    }
    required = ["pymupdf", "pypdf", "pdf2docx", "ocrmypdf", "pymupdf4llm", "openpyxl", "python-pptx"]
    required_bins = ["qpdf", "tesseract", "gs"]
    ok = all(packages[item]["ok"] for item in required) and all(binaries[item]["ok"] for item in required_bins)
    return respond(ok, engine=ENGINE_VERSION, python=sys.executable, packages=packages, binaries=binaries)


def cmd_metadata(args: argparse.Namespace) -> int:
    import fitz

    doc = fitz.open(args.input)
    return respond(True, input=args.input, pages=doc.page_count, metadata=doc.metadata)


def cmd_merge(args: argparse.Namespace) -> int:
    from pypdf import PdfReader, PdfWriter

    output = Path(args.output)
    ensure_parent(output)
    writer = PdfWriter()
    total_pages = 0

    for input_path in args.input:
        reader = PdfReader(input_path)
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                pass
        for page in reader.pages:
            writer.add_page(page)
            total_pages += 1

    with output.open("wb") as handle:
        writer.write(handle)

    return respond(True, output=str(output), inputs=args.input, pages=total_pages)


def cmd_merge_pages(args: argparse.Namespace) -> int:
    from pypdf import PdfReader, PdfWriter

    if not args.page_item:
        return respond(False, message="choose at least one page")

    output = Path(args.output)
    ensure_parent(output)
    writer = PdfWriter()
    readers: dict[str, Any] = {}
    ordered_pages: list[dict[str, Any]] = []

    for item in args.page_item:
        if "::" not in item:
            return respond(False, message=f"invalid page item: {item}")
        path_text, page_text = item.rsplit("::", 1)
        page_number = int(page_text)
        reader = readers.get(path_text)
        if reader is None:
            reader = PdfReader(path_text)
            if reader.is_encrypted:
                try:
                    reader.decrypt("")
                except Exception:
                    pass
            readers[path_text] = reader
        if page_number < 1 or page_number > len(reader.pages):
            return respond(False, message=f"page {page_number} out of range for {path_text}")
        writer.add_page(reader.pages[page_number - 1])
        ordered_pages.append({"input": path_text, "page": page_number})

    with output.open("wb") as handle:
        writer.write(handle)

    return respond(True, output=str(output), pages=len(ordered_pages), ordered_pages=ordered_pages)


def cmd_split(args: argparse.Namespace) -> int:
    import fitz

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(args.input)
    outputs: list[str] = []
    stem = safe_stem(args.input)

    for index in parse_page_spec(args.pages, doc.page_count):
        out_doc = fitz.open()
        out_doc.insert_pdf(doc, from_page=index, to_page=index)
        out_path = output_dir / f"{stem}-page-{index + 1:03d}.pdf"
        out_doc.save(str(out_path), garbage=4, deflate=True)
        out_doc.close()
        outputs.append(str(out_path))

    doc.close()
    return respond(True, output_dir=str(output_dir), outputs=outputs, pages=len(outputs))


def cmd_extract_pages(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    indexes = parse_page_spec(args.pages, doc.page_count)
    out_doc = fitz.open()
    for index in indexes:
        out_doc.insert_pdf(doc, from_page=index, to_page=index)
    out_doc.save(str(output), garbage=4, deflate=True)
    out_doc.close()
    doc.close()
    return respond(True, output=str(output), pages=[i + 1 for i in indexes])


def cmd_delete_pages(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    delete_indexes = set(parse_page_spec(args.pages, doc.page_count))
    keep_indexes = [index for index in range(doc.page_count) if index not in delete_indexes]
    if not keep_indexes:
        return respond(False, message="delete-pages would remove every page")

    out_doc = fitz.open()
    for index in keep_indexes:
        out_doc.insert_pdf(doc, from_page=index, to_page=index)
    out_doc.save(str(output), garbage=4, deflate=True)
    out_doc.close()
    doc.close()
    return respond(True, output=str(output), deleted=[i + 1 for i in sorted(delete_indexes)], pages=len(keep_indexes))


def cmd_rotate_pages(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    indexes = parse_page_spec(args.pages, doc.page_count)
    for index in indexes:
        page = doc[index]
        page.set_rotation((page.rotation + args.degrees) % 360)
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), pages=[i + 1 for i in indexes], degrees=args.degrees)


def cmd_crop_pages(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    indexes = parse_page_spec(args.pages, doc.page_count)
    for index in indexes:
        page = doc[index]
        rect = page.rect
        crop = fitz.Rect(
            rect.x0 + args.left,
            rect.y0 + args.top,
            rect.x1 - args.right,
            rect.y1 - args.bottom,
        )
        if crop.width <= 10 or crop.height <= 10:
            return respond(False, message=f"crop margins leave page {index + 1} too small")
        page.set_cropbox(crop)
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), pages=[i + 1 for i in indexes])


def find_text_block(page: Any, x: float, y: float) -> dict[str, Any] | None:
    """Locate the text block containing (or nearest to) a page point."""
    import fitz

    point = fitz.Rect(x - 1, y - 1, x + 1, y + 1)
    best = None
    best_score = None
    for block in page.get_text("dict").get("blocks", []):
        if block.get("type") != 0:
            continue
        bbox = block.get("bbox")
        if not bbox:
            continue
        # Generated PDFs often contain overlapping whitespace-only blocks;
        # only real text is clickable.
        block_text = "".join(
            span.get("text", "")
            for line in block.get("lines", [])
            for span in line.get("spans", [])
        )
        if not block_text.strip():
            continue
        rect = fitz.Rect(bbox)
        if rect.contains(fitz.Point(x, y)):
            score = -1.0
        else:
            score = rect_proximity(point, rect)
            if score > 900:  # more than ~30pt away — not what was clicked
                continue
        if best_score is None or score < best_score:
            best = block
            best_score = score
    return best


def block_payload(block: dict[str, Any]) -> dict[str, Any]:
    lines = []
    span_best = None
    span_area = 0.0
    sizes = []
    for line in block.get("lines", []):
        text = "".join(span.get("text", "") for span in line.get("spans", []))
        lines.append(text)
        for span in line.get("spans", []):
            bbox = span.get("bbox")
            if not bbox:
                continue
            area = max(0.0, (bbox[2] - bbox[0]) * (bbox[3] - bbox[1]))
            sizes.append(float(span.get("size") or 0))
            if area > span_area:
                span_area = area
                span_best = span

    raw_font = (span_best or {}).get("font") or ""
    line_count = len(lines)
    bbox = block.get("bbox")
    line_height = 0.0
    if line_count > 1 and bbox:
        line_height = (bbox[3] - bbox[1]) / line_count

    return {
        "rect": list(block.get("bbox") or ()),
        "text": "\n".join(lines),
        "fontname": raw_font,
        "fontsize": float((span_best or {}).get("size") or 12),
        "color": "#%06x" % ((span_best or {}).get("color") or 0),
        "bold": "bold" in raw_font.lower(),
        "italic": "italic" in raw_font.lower() or "oblique" in raw_font.lower(),
        "line_count": line_count,
        "line_height": line_height,
    }


def cmd_page_blocks(args: argparse.Namespace) -> int:
    """All real text blocks on a page, so the app can hit-test clicks locally."""
    import fitz

    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        doc.close()
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]
    blocks = []
    for block in page.get_text("dict").get("blocks", []):
        if block.get("type") != 0:
            continue
        text = "".join(
            span.get("text", "")
            for line in block.get("lines", [])
            for span in line.get("spans", [])
        )
        if not text.strip():
            continue
        blocks.append(block_payload(block))
    doc.close()
    return respond(True, blocks=blocks)


def cmd_block_at(args: argparse.Namespace) -> int:
    import fitz

    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        doc.close()
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]
    block = find_text_block(page, args.x, args.y)
    if not block:
        doc.close()
        return respond(True, found=False)
    payload = block_payload(block)
    doc.close()
    return respond(True, found=True, **payload)


def parse_hex_color(value: str | None) -> tuple[float, float, float] | None:
    if not value:
        return None
    raw = value.lstrip("#")
    if len(raw) != 6:
        return None
    try:
        return (int(raw[0:2], 16) / 255.0, int(raw[2:4], 16) / 255.0, int(raw[4:6], 16) / 255.0)
    except ValueError:
        return None


def block_erase_rect(rect: Any) -> Any:
    """Erase rect for a block, shrunk INWARD so redaction never removes
    neighboring glyphs that merely touch the block's bounding box (redaction
    deletes any glyph whose bbox intersects the rect)."""
    import fitz

    shrink = min(1.25, rect.height * 0.06, rect.width * 0.06)
    return fitz.Rect(rect.x0 + shrink, rect.y0 + shrink, rect.x1 - shrink, rect.y1 - shrink)


def _family_hint(raw_font: str | None) -> str | None:
    """Guess a friendly family name from a sampled PDF font, so an italic/bold
    request with no explicit --font can still select a matching face."""
    name = (raw_font or "").lower()
    if not name:
        return None
    if "times" in name or "serif" in name or name.startswith("ti"):
        return "Times New Roman"
    if "courier" in name or "mono" in name or name.startswith("co"):
        return "Courier New"
    return "Helvetica"


def resolve_block_style(doc: Any, page: Any, rect: Any, args: argparse.Namespace, text: str, scratch_files: list[str]) -> tuple[dict[str, Any], dict[str, Any]]:
    """Sample the block's style (with optional overrides) and pick a font."""
    style = style_for_rect(page, rect)
    style["underline"] = bool(getattr(args, "underline", False))
    if args.font_size:
        style["fontsize"] = args.font_size
    override_color = parse_hex_color(args.color)
    if override_color:
        style["color"] = override_color
    if args.bold and "bold" not in (style.get("rawfont") or "").lower():
        base = strip_subset_prefix(style.get("rawfont") or "Helvetica")
        style["rawfont"] = f"{base.split('-')[0]}-Bold"
        style["fontname"] = builtin_font_alias(style["rawfont"])

    wants_bold = args.bold or "bold" in (style.get("rawfont") or "").lower()
    wants_italic = bool(getattr(args, "italic", False)) \
        or "italic" in (style.get("rawfont") or "").lower() \
        or "oblique" in (style.get("rawfont") or "").lower()

    # An explicit font family wins over the sampled style: the user chose it.
    if args.font:
        font_choice = resolve_add_text_font(args.font, wants_bold, wants_italic)
        style["fontname"] = font_choice["fontname"]
        style["rawfont"] = args.font
        font_kwargs = {"fontname": font_choice["fontname"]}
        if font_choice.get("fontfile"):
            font_kwargs["fontfile"] = font_choice["fontfile"]
        return style, font_kwargs

    # Italic with no explicit family: derive a matching italic face from the
    # sampled font (the embedded font almost never has an italic variant).
    if getattr(args, "italic", False):
        family = _family_hint(style.get("rawfont")) or "Helvetica"
        font_choice = resolve_add_text_font(family, wants_bold, True)
        style["fontname"] = font_choice["fontname"]
        font_kwargs = {"fontname": font_choice["fontname"]}
        if font_choice.get("fontfile"):
            font_kwargs["fontfile"] = font_choice["fontfile"]
        return style, font_kwargs

    font_choice = resolve_replacement_font(doc, page, style, text, scratch_files)
    font_kwargs = {"fontname": font_choice["fontname"]}
    if font_choice.get("fontfile"):
        font_kwargs["fontfile"] = font_choice["fontfile"]
    return style, font_kwargs


def render_block_text(page: Any, origin_x: float, origin_y: float, width: float, text: str, style: dict[str, Any], font_kwargs: dict[str, Any], line_height: float, expand: float, align: int = 0, color_runs: list | None = None) -> bool:
    """Draw block text with its top-left at (origin_x, origin_y), wrapping at
    the block's width, growing toward the bottom of the page, shrinking only
    if it cannot fit. `align` maps to fitz.TEXT_ALIGN_*; center/right align
    within the block's width box. Draws an underline rule when the resolved
    style asks for it. When `color_runs` is set, each character is drawn in the
    color of the last run covering its global offset (else the style color)."""
    import fitz

    font_size = float(style["fontsize"])
    target = fitz.Rect(
        origin_x,
        origin_y - max(1.5, font_size * 0.1),
        min(origin_x + width + expand, page.rect.x1 - 12),
        page.rect.y1 - 12,
    )
    if color_runs:
        return render_block_text_colored(
            page, target, origin_x, origin_y, text, style, font_kwargs, line_height, align, color_runs
        )
    extra = {}
    if line_height and font_size > 0:
        extra["lineheight"] = max(0.9, min(2.5, line_height / font_size))
    size = font_size
    while size >= 6:
        result = page.insert_textbox(
            target,
            text,
            fontsize=size,
            color=style["color"],
            align=align,
            **font_kwargs,
            **extra,
        )
        if result >= 0:
            if style.get("underline"):
                bottom = min(page.rect.y1 - 2, origin_y + max(target.height, size * 1.6 * (text.count("\n") + 1)) + 4)
                draw_text_underlines(page, fitz.Rect(target.x0, target.y0, target.x1, bottom), style.get("color"))
            return True
        size -= 0.5
    return False


def color_run_rgb(color_runs: list, default_color: tuple[float, float, float]) -> list[dict[str, Any]]:
    """Normalize the incoming color-run dicts to {start, length, rgb}, dropping
    any with an unparseable hex. Order is preserved so a later run wins."""
    prepared: list[dict[str, Any]] = []
    for run in color_runs or []:
        try:
            start = int(run["start"])
            length = int(run["length"])
        except (KeyError, TypeError, ValueError):
            continue
        rgb = parse_hex_color(run.get("hex"))
        if rgb is None or length <= 0:
            continue
        prepared.append({"start": start, "length": length, "rgb": rgb})
    return prepared


def _wrap_cells(text: str, font: Any, fontsize: float, max_width: float) -> list[list[tuple[str, int]]]:
    """Greedy word-wrap `text` into visual lines, carrying each character's
    GLOBAL offset in the original string (newlines counted, not drawn) so a
    per-character color lookup stays exact across wrapped lines."""
    lines: list[list[tuple[str, int]]] = []
    current: list[tuple[str, int]] = []
    current_width = 0.0
    last_space = -1  # index in `current` of the last space, for word breaks
    i = 0
    n = len(text)
    while i < n:
        ch = text[i]
        if ch == "\n":
            lines.append(current)
            current, current_width, last_space = [], 0.0, -1
            i += 1
            continue
        w = font.text_length(ch, fontsize=fontsize)
        if current and current_width + w > max_width:
            if last_space >= 0:
                head = current[:last_space]        # drop the breaking space
                tail = current[last_space + 1:]
                lines.append(head)
                current = tail
                current_width = sum(font.text_length(c, fontsize=fontsize) for c, _ in current)
                last_space = -1
            else:
                lines.append(current)
                current, current_width, last_space = [], 0.0, -1
            continue  # re-test the same char against the fresh line
        if ch == " ":
            last_space = len(current)
        current.append((ch, i))
        current_width += w
        i += 1
    lines.append(current)
    return lines


def render_block_text_colored(page: Any, target: Any, origin_x: float, origin_y: float, text: str, style: dict[str, Any], font_kwargs: dict[str, Any], line_height: float, align: int, color_runs: list) -> bool:
    """Manual per-character-color layout: wrap to the box, then draw each line
    as contiguous same-color segments. Used only when color runs are present,
    so the fast single-color insert_textbox path is untouched otherwise."""
    import fitz

    runs = color_run_rgb(color_runs, style["color"])

    def color_at(offset: int) -> tuple[float, float, float]:
        chosen = style["color"]
        for run in runs:
            if run["start"] <= offset < run["start"] + run["length"]:
                chosen = run["rgb"]  # last covering run wins
        return chosen

    try:
        font = fitz.Font(fontfile=font_kwargs["fontfile"]) if font_kwargs.get("fontfile") else fitz.Font(font_kwargs["fontname"])
    except Exception:
        font = fitz.Font("helv")

    box_width = max(1.0, target.x1 - target.x0)
    orig_size = float(style["fontsize"])
    lh_factor = max(0.9, min(2.5, line_height / orig_size)) if (line_height and orig_size > 0) else 1.16
    ascender = font.ascender or 0.8

    size = orig_size
    while size >= 6:
        lines = _wrap_cells(text, font, size, box_width)
        line_step = size * lh_factor
        # First baseline sits one ascender below the box top; ensure the last
        # line still fits above the page's bottom margin before committing.
        first_baseline = origin_y + ascender * size
        last_baseline = first_baseline + line_step * (len(lines) - 1)
        if last_baseline <= page.rect.y1 - 6 or size <= 6:
            baseline = first_baseline
            for cells in lines:
                if cells:
                    line_w = sum(font.text_length(c, fontsize=size) for c, _ in cells)
                    if align == fitz.TEXT_ALIGN_CENTER:
                        x = origin_x + max(0.0, (box_width - line_w) / 2)
                    elif align == fitz.TEXT_ALIGN_RIGHT:
                        x = origin_x + max(0.0, box_width - line_w)
                    else:
                        x = origin_x
                    seg = ""
                    seg_color: tuple[float, float, float] | None = None
                    for ch, idx in cells:
                        col = color_at(idx)
                        if seg and col != seg_color:
                            page.insert_text((x, baseline), seg, fontsize=size, color=seg_color, **font_kwargs)
                            x += font.text_length(seg, fontsize=size)
                            seg = ""
                        seg += ch
                        seg_color = col
                    if seg and seg_color is not None:
                        page.insert_text((x, baseline), seg, fontsize=size, color=seg_color, **font_kwargs)
                baseline += line_step
            if style.get("underline"):
                draw_text_underlines(page, fitz.Rect(target.x0, target.y0, target.x1, min(page.rect.y1 - 2, baseline + 4)), style.get("color"))
            return True
        size -= 0.5
    return False


def cmd_replace_block(args: argparse.Namespace) -> int:
    """Replace an entire text block: redact its area, re-render the new text
    with the block's own style (or explicit overrides), wrapping inside the
    original box, like an in-place block editor."""
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        doc.close()
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]

    rect_values = parse_rect(args.rect)
    if rect_values is None:
        doc.close()
        return respond(False, message="--rect is required")
    rect = fitz.Rect(rect_values)

    scratch_files: list[str] = []
    new_text = args.text
    # Per-word color spans (JSON: [{"start":int,"length":int,"hex":"#rrggbb"}]).
    # Their presence forces full styled rendering (no per-line fast path).
    color_runs = None
    if getattr(args, "color_runs", None):
        try:
            parsed = json.loads(args.color_runs)
            color_runs = parsed if isinstance(parsed, list) and parsed else None
        except (ValueError, TypeError):
            color_runs = None
    style, font_kwargs = resolve_block_style(doc, page, rect, args, new_text, scratch_files)

    # Identify which lines inside the rect actually belong to this block —
    # content merely sharing the area must survive the edit untouched.
    original = args.original_text or new_text
    lines_info = block_lines_info(page, rect) or []
    block_lines, _foreign, complete = partition_block_lines(lines_info, original, rect)
    if not complete and args.original_text is None and len(lines_info) == len(new_text.split("\n")):
        # No original text supplied: fall back to count-based mapping.
        block_lines, complete = lines_info, True

    # When the edited text has the same number of lines as the block, render
    # line-by-line at each line's own origin and style — indents and
    # right-aligned amounts stay exactly where they were.
    new_lines = new_text.split("\n")
    per_line = (
        complete
        and not args.font_size and not args.color and not args.bold and not args.font
        and not getattr(args, "italic", False) and not getattr(args, "underline", False)
        and getattr(args, "align", "left") == "left"
        and not color_runs
        and len(new_lines) == len(block_lines)
        and all(line.get("origin") for line in block_lines)
    )
    per_line_fonts: list[dict[str, Any]] = []
    if per_line:
        font_cache: dict[str, dict[str, Any]] = {}
        for line, text_line in zip(block_lines, new_lines):
            raw = line["rawfont"] or ""
            if raw not in font_cache:
                line_style = {
                    "fontname": builtin_font_alias(raw),
                    "fontsize": line["fontsize"],
                    "color": line["color"],
                    "rawfont": raw,
                }
                font_cache[raw] = resolve_replacement_font(doc, page, line_style, new_text, scratch_files)
            per_line_fonts.append(font_cache[raw])
        # An edited line that would spill past the page edge needs wrapping —
        # fall back to flowed rendering in that case.
        for line, text_line, font_choice in zip(block_lines, new_lines, per_line_fonts):
            length = replacement_text_length(text_line, font_choice, line["fontsize"])
            if line["origin"][0] + length > page.rect.x1 - 8:
                per_line = False
                break

    if block_lines:
        erase_lines(page, block_lines)
    else:
        page.add_redact_annot(block_erase_rect(rect), fill=(1, 1, 1))
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)

    # Shade behind the block (drawn AFTER erasing, BEFORE text) so text
    # color and background persist together in one operation.
    bg = parse_hex_color(args.background) if getattr(args, "background", None) else None
    if bg is not None:
        page.draw_rect(fitz.Rect(rect) + (-2, -2, 2, 2), fill=bg, color=None, overlay=True)

    inserted = False
    if per_line and new_text.strip():
        inserted = True
        for line, text_line, font_choice in zip(block_lines, new_lines, per_line_fonts):
            if not text_line.strip():
                continue
            kwargs = {"fontname": font_choice["fontname"]}
            if font_choice.get("fontfile"):
                kwargs["fontfile"] = font_choice["fontfile"]
            page.insert_text(
                (line["origin"][0], line["origin"][1]),
                text_line,
                fontsize=line["fontsize"],
                color=line["color"],
                **kwargs,
            )
    elif new_text.strip():
        inserted = render_block_text(
            page, rect.x0, rect.y0, rect.width, new_text,
            style, font_kwargs, args.line_height, args.expand,
            align=align_constant(getattr(args, "align", "left")),
            color_runs=color_runs,
        )

    # Track Changes: leave a dashed blue review box whose note holds the
    # original wording, so every edit stays reviewable.
    if inserted and args.track_original:
        grown = fitz.Rect(rect) + (-2, -2, 2, 2)
        annot = page.add_rect_annot(grown)
        annot.set_colors(stroke=(0.15, 0.4, 0.85))
        annot.set_border(width=1, dashes=[3, 2])
        original = args.track_original
        snippet = original if len(original) <= 600 else original[:600] + "…"
        annot.set_info(title="Tracked change", content=f"Original text:\n{snippet}")
        annot.update()

    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    for path in scratch_files:
        try:
            os.unlink(path)
        except OSError:
            pass
    if new_text.strip() and not inserted:
        return respond(False, message="replacement text did not fit")
    return respond(True, output=str(output), inserted=inserted)


def block_lines_info(page: Any, rect: Any) -> list[dict[str, Any]] | None:
    """Per-line text, baseline origin, and dominant style for every text line
    inside rect — the same content the block erase removes — so edits and
    moves can reproduce each line's exact position and styling."""
    import fitz

    target = fitz.Rect(rect) + (-1, -1, 1, 1)
    lines = []
    for block in page.get_text("dict").get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = [s for s in line.get("spans", []) if s.get("text")]
            if not spans:
                continue
            bbox = line.get("bbox")
            if not bbox:
                continue
            center = fitz.Point((bbox[0] + bbox[2]) / 2, (bbox[1] + bbox[3]) / 2)
            if not target.contains(center):
                continue
            dominant = max(
                spans,
                key=lambda s: max(0.0, (s["bbox"][2] - s["bbox"][0]) * (s["bbox"][3] - s["bbox"][1])),
            )
            lines.append({
                "text": "".join(s["text"] for s in spans),
                "origin": spans[0].get("origin"),
                "bbox": bbox,
                "rawfont": dominant.get("font"),
                "fontsize": float(dominant.get("size") or 11),
                "color": pdf_color_to_rgb(dominant.get("color")),
                "sort_key": (round(bbox[1], 1), round(bbox[0], 1)),
            })
    lines.sort(key=lambda line: line["sort_key"])
    return lines or None


def partition_block_lines(lines_info: list[dict[str, Any]], original_text: str, rect: Any = None) -> tuple[list[dict[str, Any]], list[dict[str, Any]], bool]:
    """Split the lines inside a rect into the block's own lines (matched
    against its original text) and FOREIGN lines that merely share the area.

    Documents can contain IDENTICAL duplicate lines (especially ones damaged
    by earlier edits) — matching is by spatial affinity: each original line
    claims the unclaimed candidate with equal text that sits closest to the
    block rect, and candidates further than 40pt are never claimed. Foreign
    content must never be erased by a move or edit."""
    import fitz

    target = fitz.Rect(rect) if rect is not None else None

    def distance(line: dict[str, Any]) -> float:
        if target is None:
            return 0.0
        bbox = fitz.Rect(line["bbox"])
        if target.intersects(bbox):
            return 0.0
        dx = max(target.x0 - bbox.x1, bbox.x0 - target.x1, 0)
        dy = max(target.y0 - bbox.y1, bbox.y0 - target.y1, 0)
        return (dx * dx + dy * dy) ** 0.5

    claimed = [False] * len(lines_info)
    matched_count = 0
    for wanted in original_text.split("\n"):
        best_index = -1
        best_distance = 40.0  # never claim text further than this
        for index, line in enumerate(lines_info):
            if claimed[index] or line["text"] != wanted:
                continue
            d = distance(line)
            if d < best_distance:
                best_distance = d
                best_index = index
        if best_index >= 0:
            claimed[best_index] = True
            matched_count += 1

    block_lines = [line for index, line in enumerate(lines_info) if claimed[index]]
    foreign = [line for index, line in enumerate(lines_info) if not claimed[index]]
    complete = matched_count == len(original_text.split("\n")) and bool(block_lines)
    return block_lines, foreign, complete


def erase_lines(page: Any, lines: list[dict[str, Any]]) -> None:
    """Redact ONLY the given lines' own glyph areas (slightly shrunk), so
    neighbors and content sharing the block's bounding box survive."""
    import fitz

    for line in lines:
        rect = fitz.Rect(line["bbox"])
        shrink = min(0.75, rect.height * 0.15)
        page.add_redact_annot(
            fitz.Rect(rect.x0 + shrink, rect.y0 + shrink, rect.x1 - shrink, rect.y1 - shrink),
            fill=(1, 1, 1),
        )
    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)


def cmd_move_block(args: argparse.Namespace) -> int:
    """Move a text block: erase ONLY the block's own glyphs (inward-shrunk
    redact), then re-render its text with matched style exactly at the drop
    point. Never copies or erases neighboring content the way a rectangle
    move would. Pure moves (text unchanged) keep each line's own font, size,
    color, and relative position."""
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        doc.close()
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]

    rect_values = parse_rect(args.rect)
    if rect_values is None:
        doc.close()
        return respond(False, message="--rect is required")
    rect = fitz.Rect(rect_values)

    text = args.text
    if not text.strip():
        doc.close()
        return respond(False, message="--text is required for move-block")

    scratch_files: list[str] = []
    original = args.original_text or text
    lines_info = block_lines_info(page, rect) or []
    block_lines, _foreign, complete = partition_block_lines(lines_info, original, rect)

    pure_move = (
        complete
        and text == original
        and not args.font_size and not args.color and not args.bold and not args.font
        and not getattr(args, "italic", False) and not getattr(args, "underline", False)
        and getattr(args, "align", "left") == "left"
        and all(line.get("origin") for line in block_lines)
    )

    style, font_kwargs = resolve_block_style(doc, page, rect, args, text, scratch_files)
    per_line_fonts: list[dict[str, Any]] = []
    if pure_move:
        font_cache: dict[str, dict[str, Any]] = {}
        for line in block_lines:
            raw = line["rawfont"] or ""
            if raw not in font_cache:
                line_style = {
                    "fontname": builtin_font_alias(raw),
                    "fontsize": line["fontsize"],
                    "color": line["color"],
                    "rawfont": raw,
                }
                font_cache[raw] = resolve_replacement_font(doc, page, line_style, line["text"], scratch_files)
            per_line_fonts.append(font_cache[raw])

    # Erase ONLY the block's own lines: content that merely shares the
    # block's bounding box (or sits in the drag path) is never touched.
    if block_lines:
        erase_lines(page, block_lines)
    else:
        page.add_redact_annot(block_erase_rect(rect), fill=(1, 1, 1))
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)

    dest_x = max(page.rect.x0 + 6, min(args.dest_x, page.rect.x1 - rect.width - 6))
    dest_y = max(page.rect.y0 + 6, min(args.dest_y, page.rect.y1 - 18))

    if pure_move:
        # Translate every line to the destination, preserving per-line style
        # and exact intra-block layout (indents, baselines).
        inserted = True
        for line, font_choice in zip(block_lines, per_line_fonts):
            kwargs = {"fontname": font_choice["fontname"]}
            if font_choice.get("fontfile"):
                kwargs["fontfile"] = font_choice["fontfile"]
            origin = line["origin"]
            page.insert_text(
                (dest_x + (origin[0] - rect.x0), dest_y + (origin[1] - rect.y0)),
                line["text"],
                fontsize=line["fontsize"],
                color=line["color"],
                **kwargs,
            )
    else:
        inserted = render_block_text(
            page, dest_x, dest_y, rect.width, text,
            style, font_kwargs, args.line_height, args.expand,
            align=align_constant(getattr(args, "align", "left")),
        )

    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    for path in scratch_files:
        try:
            os.unlink(path)
        except OSError:
            pass
    if not inserted:
        return respond(False, message="moved text did not fit at the destination")
    return respond(True, output=str(output), dest=[dest_x, dest_y])


def cmd_add_symbol(args: argparse.Namespace) -> int:
    """Fill & Sign stamps: checkmark, cross, or dot at a click point, drawn
    with the base-14 ZapfDingbats font so they render everywhere."""
    import fitz

    glyphs = {"check": "3", "cross": "7", "dot": "l"}
    glyph = glyphs.get(args.kind)
    if glyph is None:
        return respond(False, message=f"kind must be one of {sorted(glyphs)}")

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        doc.close()
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]
    color = parse_hex_color(args.color) or (0.05, 0.05, 0.05)
    page.insert_text(
        (args.x, args.y + args.size * 0.72),
        glyph,
        fontname="zadb",
        fontsize=args.size,
        color=color,
    )
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), kind=args.kind)


def cmd_compress(args: argparse.Namespace) -> int:
    """Reduce file size with Ghostscript. Quality: small (screen 72dpi),
    medium (ebook 150dpi), high (printer 300dpi)."""
    settings = {"small": "/screen", "medium": "/ebook", "high": "/printer"}
    preset = settings.get(args.quality)
    if preset is None:
        return respond(False, message=f"quality must be one of {sorted(settings)}")

    gs = shutil.which("gs") or "/opt/homebrew/bin/gs"
    if not os.path.exists(gs):
        return respond(False, message="Ghostscript (gs) not found; brew install ghostscript")

    output = Path(args.output)
    ensure_parent(output)
    before = os.path.getsize(args.input)
    result = subprocess.run(
        [
            gs, "-q", "-dNOPAUSE", "-dBATCH", "-dSAFER",
            "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5",
            f"-dPDFSETTINGS={preset}",
            f"-sOutputFile={output}",
            str(args.input),
        ],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0 or not output.exists():
        return respond(False, message=f"Ghostscript failed: {result.stderr[:400]}")
    after = os.path.getsize(output)
    if after >= before:
        # Compression did not help (already optimized) — keep the original.
        shutil.copyfile(args.input, output)
        after = before
    return respond(True, output=str(output), before_bytes=before, after_bytes=after)


def cmd_add_page_numbers(args: argparse.Namespace) -> int:
    import fitz

    positions = {"bottom-center", "bottom-left", "bottom-right", "top-center", "top-left", "top-right"}
    if args.position not in positions:
        return respond(False, message=f"position must be one of {sorted(positions)}")
    formats = {"n", "page-n", "n-of-total"}
    if args.number_format not in formats:
        return respond(False, message=f"format must be one of {sorted(formats)}")

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    total = doc.page_count
    size = args.font_size
    margin = args.margin
    for index in range(total):
        page = doc[index]
        number = args.start + index
        if args.number_format == "page-n":
            label = f"Page {number}"
        elif args.number_format == "n-of-total":
            label = f"{number} of {args.start + total - 1}"
        else:
            label = str(number)
        width = fitz.get_text_length(label, fontname="helv", fontsize=size)
        rect = page.rect
        if "left" in args.position:
            x = rect.x0 + margin
        elif "right" in args.position:
            x = rect.x1 - margin - width
        else:
            x = rect.x0 + (rect.width - width) / 2
        y = rect.y0 + margin if args.position.startswith("top") else rect.y1 - margin + size * 0.72 - size
        page.insert_text((x, y + size * 0.72), label, fontname="helv", fontsize=size, color=(0.2, 0.2, 0.2))
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), pages=total)


def cmd_resize_pages(args: argparse.Namespace) -> int:
    """Rescale every page onto a uniform target size (content aspect-fit,
    centered) — makes merged PDFs of mixed page sizes one consistent size."""
    import fitz

    if args.width <= 10 or args.height <= 10:
        return respond(False, message="width/height must be in points and > 10")

    output = Path(args.output)
    ensure_parent(output)
    src = fitz.open(args.input)
    out = fitz.open()
    for index in range(src.page_count):
        page = src[index]
        target = out.new_page(width=args.width, height=args.height)
        src_rect = page.rect
        if src_rect.width <= 0 or src_rect.height <= 0:
            continue
        scale = min(args.width / src_rect.width, args.height / src_rect.height)
        drawn_w = src_rect.width * scale
        drawn_h = src_rect.height * scale
        x0 = (args.width - drawn_w) / 2
        y0 = (args.height - drawn_h) / 2
        target.show_pdf_page(fitz.Rect(x0, y0, x0 + drawn_w, y0 + drawn_h), src, index)
    out.save(str(output), garbage=4, deflate=True)
    pages = out.page_count
    out.close()
    src.close()
    return respond(True, output=str(output), pages=pages, width=args.width, height=args.height)


def cmd_block_background(args: argparse.Namespace) -> int:
    """Fill a colored rectangle BEHIND existing content — background shading
    for a block, separate from text color."""
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    color = parse_hex_color(args.color)
    if color is None:
        return respond(False, message="--color must be #rrggbb")
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        doc.close()
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    rect_values = parse_rect(args.rect)
    if rect_values is None:
        doc.close()
        return respond(False, message="--rect is required")
    page = doc[args.page - 1]
    rect = fitz.Rect(rect_values) + (-args.padding, -args.padding, args.padding, args.padding)
    page.draw_rect(rect, fill=color, color=None, overlay=False)
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output))


def cmd_set_password(args: argparse.Namespace) -> int:
    """Protect the PDF with a password (AES-256)."""
    import fitz

    if not args.password:
        return respond(False, message="--password cannot be empty")
    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    permissions = int(
        fitz.PDF_PERM_ACCESSIBILITY
        | fitz.PDF_PERM_PRINT
        | fitz.PDF_PERM_COPY
        | fitz.PDF_PERM_ANNOTATE
    )
    doc.save(
        str(output),
        encryption=fitz.PDF_ENCRYPT_AES_256,
        owner_pw=args.password,
        user_pw=args.password,
        permissions=permissions,
        garbage=4,
        deflate=True,
    )
    doc.close()
    return respond(True, output=str(output))


def cmd_redline(args: argparse.Namespace) -> int:
    """Review marks: squiggly underline, insert caret, or replace (strikeout +
    caret + note) over the given line rects."""
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        doc.close()
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]

    rects = []
    for part in (args.rects or "").split(";"):
        values = parse_rect(part) if part.strip() else None
        if values:
            rects.append(fitz.Rect(values))

    red = (0.78, 0.16, 0.16)
    made = 0
    if args.kind == "strikeout":
        for rect in rects:
            annot = page.add_strikeout_annot(rect)
            annot.set_colors(stroke=red)
            annot.update()
            made += 1
    elif args.kind == "squiggly":
        for rect in rects:
            annot = page.add_squiggly_annot(rect)
            annot.set_colors(stroke=red)
            annot.update()
            made += 1
    elif args.kind == "caret":
        point = fitz.Point(args.x, args.y)
        annot = page.add_caret_annot(point)
        annot.set_colors(stroke=(0.1, 0.3, 0.8))
        if args.note:
            annot.set_info(content=args.note)
        annot.update()
        made = 1
    elif args.kind == "replace":
        if not rects:
            doc.close()
            return respond(False, message="--rects required for replace")
        for rect in rects:
            annot = page.add_strikeout_annot(rect)
            annot.set_colors(stroke=red)
            if args.note:
                annot.set_info(content=f"Replace with: {args.note}")
            annot.update()
            made += 1
        last = rects[-1]
        caret = page.add_caret_annot(fitz.Point(last.x1, last.y1))
        caret.set_colors(stroke=(0.1, 0.3, 0.8))
        if args.note:
            caret.set_info(content=f"Replace with: {args.note}")
        caret.update()
    else:
        doc.close()
        return respond(False, message="kind must be strikeout, squiggly, caret, or replace")

    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), kind=args.kind, marks=made)


def cmd_replace_text(args: argparse.Namespace) -> int:
    import fitz

    if not args.find:
        return respond(False, message="--find cannot be empty")

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    matches = 0
    inserted = 0
    rect_values = parse_rect(args.rect)
    target_indexes = [args.page - 1] if args.page else list(range(doc.page_count))
    scratch_files: list[str] = []

    for page_index in target_indexes:
        if page_index < 0 or page_index >= doc.page_count:
            return respond(False, message=f"page must be between 1 and {doc.page_count}")
        page = doc[page_index]
        if rect_values:
            click_rect = fitz.Rect(rect_values)
            found_rects = [nearest_search_rect(page, click_rect, args.find) or click_rect]
        else:
            found_rects = search_rects(page, args.find)
        if not found_rects:
            continue

        matches += len(found_rects)
        styles = [style_for_rect(page, rect) for rect in found_rects]
        font_choices = [
            resolve_replacement_font(doc, page, style, args.replace, scratch_files)
            if args.match_style
            else {"fontname": "helv", "fontfile": None}
            for style in styles
        ]

        for rect in found_rects:
            # Exact hit rect: outward padding would delete neighboring glyphs
            # that merely touch it (redaction removes any intersecting glyph).
            page.add_redact_annot(rect, fill=(1, 1, 1))

        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)

        for rect, style, font_choice in zip(found_rects, styles, font_choices):
            base_font_size = float(style["fontsize"])
            color = style["color"] if args.match_style else (0, 0, 0)
            font_kwargs = {"fontname": font_choice["fontname"]}
            if font_choice.get("fontfile"):
                font_kwargs["fontfile"] = font_choice["fontfile"]

            font_size = base_font_size if args.match_style else args.font_size

            # Single-line replacements sit on the original baseline at the
            # original size; shrink only if the new text
            # would run off the right edge of the page.
            single_line = "\n" not in args.replace
            origin = style.get("origin")
            if single_line:
                if args.auto_size:
                    max_width = page.rect.x1 - 12 - rect.x0
                    while font_size > 6 and replacement_text_length(args.replace, font_choice, font_size) > max_width:
                        font_size -= 0.5
                baseline_y = origin[1] if origin else rect.y1 - max(1.5, base_font_size * 0.18)
                page.insert_text(
                    (rect.x0, baseline_y),
                    args.replace,
                    fontsize=max(font_size, 6),
                    color=color,
                    **font_kwargs,
                )
                inserted += 1
                continue

            target_width = max(
                rect.width + args.expand,
                replacement_text_length(args.replace, font_choice, base_font_size) + 6,
                80,
            )
            target = fitz.Rect(
                rect.x0,
                max(page.rect.y0 + 2, rect.y0 - max(2, base_font_size * 0.12)),
                min(page.rect.x1 - 12, rect.x0 + target_width),
                min(page.rect.y1 - 2, rect.y1 + max(5, base_font_size * 0.35)),
            )

            result = -1
            while args.auto_size and font_size >= 6:
                result = page.insert_textbox(
                    target,
                    args.replace,
                    fontsize=font_size,
                    color=color,
                    align=fitz.TEXT_ALIGN_LEFT,
                    **font_kwargs,
                )
                if result >= 0:
                    break
                font_size -= 0.5

            if not args.auto_size:
                result = page.insert_textbox(
                    target,
                    args.replace,
                    fontsize=font_size,
                    color=color,
                    align=fitz.TEXT_ALIGN_LEFT,
                    **font_kwargs,
                )

            if result < 0:
                page.insert_text(
                    (target.x0, target.y1 - 3),
                    args.replace,
                    fontsize=max(font_size, 6),
                    color=color,
                    **font_kwargs,
                )
            inserted += 1

    if matches == 0:
        doc.close()
        for path in scratch_files:
            try:
                os.unlink(path)
            except OSError:
                pass
        return respond(False, message=f"No matches found for {args.find!r}")

    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    for path in scratch_files:
        try:
            os.unlink(path)
        except OSError:
            pass
    return respond(True, output=str(output), matches=matches, inserted=inserted)


def cmd_redact_text(args: argparse.Namespace) -> int:
    import fitz

    if not args.find:
        return respond(False, message="--find cannot be empty")

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    rect_values = parse_rect(args.rect)
    indexes = [args.page - 1] if args.page else parse_page_spec(args.pages, doc.page_count)
    matches = 0

    for index in indexes:
        if index < 0 or index >= doc.page_count:
            return respond(False, message=f"page must be between 1 and {doc.page_count}")
        page = doc[index]
        if rect_values:
            click_rect = fitz.Rect(rect_values)
            rects = [nearest_search_rect(page, click_rect, args.find) or click_rect]
        else:
            rects = search_rects(page, args.find)
        matches += len(rects)
        for rect in rects:
            padded = fitz.Rect(rect.x0 - 1, rect.y0 - 1, rect.x1 + 1, rect.y1 + 1)
            page.add_redact_annot(padded, text=args.label or None, fill=(1, 1, 1), text_color=(0, 0, 0))
        if rects:
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)

    if matches == 0:
        doc.close()
        return respond(False, message=f"No matches found for {args.find!r}")

    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), matches=matches, pages=[i + 1 for i in indexes])


# Friendly font names → base-14 aliases by variant
# (regular, bold, italic, bold-italic). Verified against fitz.Font():
# Helvetica-Oblique is "heit" (not "heli"), bold-italic "hebi".
_ADD_TEXT_BUILTIN = {
    "helvetica": ("helv", "hebo", "heit", "hebi"),
    "times new roman": ("tiro", "tibo", "tiit", "tibi"),
    "times": ("tiro", "tibo", "tiit", "tibi"),
    "courier new": ("cour", "cobo", "coit", "cobi"),
    "courier": ("cour", "cobo", "coit", "cobi"),
}

# Folders that hold installed macOS font files, richest first.
_ADD_TEXT_FONT_DIRS = (
    "/System/Library/Fonts/Supplemental",
    "/System/Library/Fonts",
    "/Library/Fonts",
    os.path.expanduser("~/Library/Fonts"),
)


def _variant_score(rest: str, bold: bool, italic: bool) -> int | None:
    """Rank a font file's trailing variant tokens against the wanted style.
    `rest` is the file stem with the family prefix already stripped (spaces
    and hyphens removed, lower-cased). Lower score = better match; None means
    the file carries a style we explicitly did NOT ask for."""
    has_bold = "bold" in rest or rest.endswith("bd") or "black" in rest or "heavy" in rest
    has_italic = "italic" in rest or "oblique" in rest or rest.endswith("it")
    if bold and italic:
        if has_bold and has_italic:
            return 0
    elif bold:
        if has_bold and not has_italic:
            return 0
        if has_bold:
            return 2
    elif italic:
        if has_italic and not has_bold:
            return 0
        if has_italic:
            return 2
    else:
        if not has_bold and not has_italic:
            return 0
    # Regular file is always an acceptable fallback for any request.
    if not has_bold and not has_italic:
        return 5
    return None


def find_family_font_file(family: str | None, bold: bool, italic: bool) -> str | None:
    """Locate an installed .ttf/.otf whose name matches `family`, honoring the
    bold/italic variant. Matches by compact (spaces/hyphens removed, case-
    insensitive) family prefix so "Times New Roman" finds "TimesNewRomanPSMT"
    etc. Returns the best-scoring file path or None."""
    compact = re.sub(r"[\s\-]", "", (family or "")).lower()
    if not compact:
        return None
    best_path: str | None = None
    best_score = 1_000
    for directory in _ADD_TEXT_FONT_DIRS:
        if not os.path.isdir(directory):
            continue
        try:
            entries = os.listdir(directory)
        except OSError:
            continue
        for entry in entries:
            stem, ext = os.path.splitext(entry)
            if ext.lower() not in (".ttf", ".otf"):
                continue
            compact_stem = re.sub(r"[\s\-_]", "", stem).lower()
            if not compact_stem.startswith(compact):
                continue
            rest = compact_stem[len(compact):]
            score = _variant_score(rest, bold, italic)
            if score is None:
                continue
            # Shorter remainder (closer family match) breaks score ties, so
            # "Arial" beats "Arial Narrow" and "Arial Bold" beats the narrow
            # variants for the same requested style.
            score = score * 100 + min(len(rest), 99)
            if score < best_score:
                best_score = score
                best_path = os.path.join(directory, entry)
    return best_path


def resolve_add_text_font(name: str | None, bold: bool, italic: bool = False) -> dict[str, Any]:
    """Map ANY font family name to insert_text kwargs.

    Resolution order: (a) base-14 builtins (Helvetica/Times/Courier) with the
    right bold/italic variant; (b) an installed TTF/OTF for the family (any
    system family — Georgia, Verdana, Arial, …) honoring the variant, falling
    back to the family's regular file; (c) the Helvetica base-14 variant."""
    key = (name or "").strip().lower()
    builtin = _ADD_TEXT_BUILTIN.get(key)
    if builtin:
        index = (2 if italic else 0) + (1 if bold else 0)
        return {"fontname": builtin[index], "fontfile": None}

    path = find_family_font_file(name, bold, italic)
    if path is None and (bold or italic):
        path = find_family_font_file(name, False, False)
    if path:
        label = re.sub(r"[^A-Za-z0-9]", "", name or "")[:16] or "F0"
        return {"fontname": label, "fontfile": path}

    helv = _ADD_TEXT_BUILTIN["helvetica"]
    index = (2 if italic else 0) + (1 if bold else 0)
    return {"fontname": helv[index], "fontfile": None}


def cmd_add_text(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        return respond(False, message=f"page must be between 1 and {doc.page_count}")

    color = parse_hex_color(args.color) or (0, 0, 0)
    font_choice = resolve_add_text_font(args.font, args.bold, getattr(args, "italic", False))
    kwargs = {"fontname": font_choice["fontname"]}
    if font_choice.get("fontfile"):
        kwargs["fontfile"] = font_choice["fontfile"]
    # --y is the TOP of the first line; PyMuPDF draws from the baseline.
    baseline_y = args.y + args.font_size * 0.83

    page = doc[args.page - 1]
    align = align_constant(getattr(args, "align", "left"))
    if align != fitz.TEXT_ALIGN_LEFT:
        # A box from the click x to the right margin gives center/right text
        # room to lay out; left keeps the exact baseline placement below
        # (insert_text already handles left-aligned multi-line text).
        box = fitz.Rect(args.x, args.y - max(1.5, args.font_size * 0.1), page.rect.x1 - 12, page.rect.y1 - 12)
        page.insert_textbox(
            box, args.text,
            fontsize=args.font_size, color=color, align=align, **kwargs,
        )
    else:
        page.insert_text(
            (args.x, baseline_y),
            args.text,
            fontsize=args.font_size,
            color=color,
            **kwargs,
        )
    if getattr(args, "underline", False):
        clip = fitz.Rect(
            args.x - 1, args.y - 2, page.rect.x1 - 2,
            args.y + args.font_size * 1.6 * (args.text.count("\n") + 1) + 4,
        )
        draw_text_underlines(page, clip, color)
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), page=args.page)


def cmd_annotate_text(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    indexes = parse_page_spec(args.pages, doc.page_count)
    matches = 0

    for index in indexes:
        page = doc[index]
        rects = search_rects(page, args.find)
        matches += len(rects)
        for rect in rects:
            if args.kind == "highlight":
                annot = page.add_highlight_annot(rect)
            elif args.kind == "underline":
                annot = page.add_underline_annot(rect)
            else:
                annot = page.add_strikeout_annot(rect)
            annot.set_info(content=args.note or f"{args.kind}: {args.find}")
            annot.update()

    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), kind=args.kind, matches=matches)


def cmd_add_note(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]
    annot = page.add_text_annot((args.x, args.y), args.text)
    annot.set_info(title="Sam PDF Studio", content=args.text)
    annot.update()
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), page=args.page)


def cmd_add_signature(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]
    rect = fitz.Rect(args.x, args.y, args.x + args.width, args.y + args.height)
    page.insert_textbox(
        rect,
        args.text,
        fontname="tiro",
        fontsize=args.font_size,
        color=(0, 0, 0),
        align=fitz.TEXT_ALIGN_CENTER,
    )
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), page=args.page)


def cmd_add_image(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]
    rect = fitz.Rect(args.x, args.y, args.x + args.width, args.y + args.height)
    page.insert_image(rect, filename=args.image)
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), page=args.page, image=args.image)


def cmd_paste_region(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    source_doc = fitz.open(args.source)
    dest_doc = fitz.open(args.input)

    if args.source_page < 1 or args.source_page > source_doc.page_count:
        source_doc.close()
        dest_doc.close()
        return respond(False, message=f"source page must be between 1 and {source_doc.page_count}")
    if args.destination_page < 1 or args.destination_page > dest_doc.page_count:
        source_doc.close()
        dest_doc.close()
        return respond(False, message=f"destination page must be between 1 and {dest_doc.page_count}")

    source_rect_values = parse_rect(args.source_rect)
    if source_rect_values is None:
        source_doc.close()
        dest_doc.close()
        return respond(False, message="--source-rect is required")

    source_page = source_doc[args.source_page - 1]
    source_rect = fitz.Rect(source_rect_values)
    clip = source_rect & source_page.rect
    if clip.width <= 1 or clip.height <= 1:
        source_doc.close()
        dest_doc.close()
        return respond(False, message="source region is outside the page")

    # Build a single-page scratch doc containing ONLY the clipped region as
    # live vector content (text stays selectable/searchable/redactable and
    # razor sharp at any zoom). Everything outside the clip
    # is redacted away so the embedded form holds no hidden content.
    region_doc = fitz.open()
    region_doc.insert_pdf(source_doc, from_page=args.source_page - 1, to_page=args.source_page - 1)
    region_page = region_doc[0]
    page_rect = region_page.rect
    outside = [
        fitz.Rect(page_rect.x0, page_rect.y0, page_rect.x1, clip.y0),
        fitz.Rect(page_rect.x0, clip.y1, page_rect.x1, page_rect.y1),
        fitz.Rect(page_rect.x0, clip.y0, clip.x0, clip.y1),
        fitz.Rect(clip.x1, clip.y0, page_rect.x1, clip.y1),
    ]
    for strip in outside:
        if strip.width > 0.5 and strip.height > 0.5:
            region_page.add_redact_annot(strip, fill=False)
    region_page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)

    same_file = Path(args.source).resolve() == Path(args.input).resolve()
    if args.erase_source:
        if not same_file:
            source_doc.close()
            dest_doc.close()
            region_doc.close()
            return respond(False, message="--erase-source can only be used when source and input are the same PDF")
        erase_page = dest_doc[args.source_page - 1]
        padded = fitz.Rect(clip.x0 - 1, clip.y0 - 1, clip.x1 + 1, clip.y1 + 1)
        erase_page.add_redact_annot(padded, fill=(1, 1, 1))
        erase_page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)

    dest_page = dest_doc[args.destination_page - 1]
    dest_x = max(dest_page.rect.x0, min(args.destination_x, dest_page.rect.x1 - clip.width))
    dest_y = max(dest_page.rect.y0, min(args.destination_y, dest_page.rect.y1 - clip.height))
    dest_rect = fitz.Rect(dest_x, dest_y, dest_x + clip.width, dest_y + clip.height)
    try:
        dest_page.show_pdf_page(dest_rect, region_doc, 0, clip=clip)
    except Exception:
        # Fallback: rasterize the region (previous behavior).
        matrix = fitz.Matrix(args.dpi / 72, args.dpi / 72)
        pixmap = source_page.get_pixmap(matrix=matrix, clip=clip, alpha=False)
        dest_page.insert_image(dest_rect, stream=pixmap.tobytes("png"))
    region_doc.close()

    dest_doc.save(str(output), garbage=4, deflate=True)
    source_doc.close()
    dest_doc.close()
    return respond(
        True,
        output=str(output),
        source_page=args.source_page,
        destination_page=args.destination_page,
        erased=args.erase_source,
        rect=[dest_rect.x0, dest_rect.y0, dest_rect.x1, dest_rect.y1],
    )


def cmd_add_link(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    if args.page < 1 or args.page > doc.page_count:
        return respond(False, message=f"page must be between 1 and {doc.page_count}")
    page = doc[args.page - 1]
    rect = fitz.Rect(args.x, args.y, args.x + args.width, args.y + args.height)
    page.insert_link({"kind": fitz.LINK_URI, "from": rect, "uri": args.url})
    doc.save(str(output), garbage=4, deflate=True)
    doc.close()
    return respond(True, output=str(output), page=args.page, url=args.url)


def run_ocrmypdf(input_path: Path, output_path: Path, language: str, force: bool) -> subprocess.CompletedProcess[str]:
    command = [
        sys.executable,
        "-m",
        "ocrmypdf",
        "--output-type",
        "pdf",
        "--deskew",
        "--rotate-pages",
        "--language",
        language,
    ]
    command.append("--force-ocr" if force else "--skip-text")
    command.extend([str(input_path), str(output_path)])

    env = os.environ.copy()
    env["PATH"] = ":".join([
        str(Path(sys.executable).resolve().parent),
        "/opt/homebrew/bin",
        "/usr/local/bin",
        env.get("PATH", ""),
    ])
    font_config = Path("/opt/homebrew/etc/fonts/fonts.conf")
    if font_config.exists():
        env.setdefault("FONTCONFIG_FILE", str(font_config))
    return subprocess.run(command, text=True, capture_output=True, env=env)


def cmd_ocr(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)

    if args.page:
        source = fitz.open(args.input)
        if args.page < 1 or args.page > source.page_count:
            return respond(False, message=f"page must be between 1 and {source.page_count}")
        with tempfile.TemporaryDirectory(prefix="sampdf-ocr-") as tmp_dir:
            single_path = Path(tmp_dir) / "page.pdf"
            ocr_path = Path(tmp_dir) / "page-ocr.pdf"
            single = fitz.open()
            single.insert_pdf(source, from_page=args.page - 1, to_page=args.page - 1)
            single.save(str(single_path), garbage=4, deflate=True)
            single.close()
            result = run_ocrmypdf(single_path, ocr_path, args.language, True)
            if result.returncode != 0:
                source.close()
                return respond(False, output=str(output), message="OCR failed", stdout=result.stdout[-4000:], stderr=result.stderr[-4000:])

            ocr_doc = fitz.open(str(ocr_path))
            final = fitz.open()
            for index in range(source.page_count):
                if index == args.page - 1:
                    final.insert_pdf(ocr_doc)
                else:
                    final.insert_pdf(source, from_page=index, to_page=index)
            final.save(str(output), garbage=4, deflate=True)
            final.close()
            ocr_doc.close()
            source.close()
        return respond(True, output=str(output), page=args.page, stdout=result.stdout[-1000:], stderr=result.stderr[-1000:])

    result = run_ocrmypdf(Path(args.input), output, args.language, args.force)
    if result.returncode != 0:
        return respond(False, output=str(output), message="OCR failed", stdout=result.stdout[-4000:], stderr=result.stderr[-4000:])
    return respond(True, output=str(output), stdout=result.stdout[-1000:], stderr=result.stderr[-1000:])


def cmd_export_images(args: argparse.Namespace) -> int:
    import fitz

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(args.input)
    zoom = args.dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    outputs: list[str] = []
    stem = safe_stem(args.input)

    for index, page in enumerate(doc, start=1):
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        suffix = "jpg" if args.format == "jpg" else "png"
        out_path = output_dir / f"{stem}-page-{index:03d}.{suffix}"
        if args.format == "jpg":
            pixmap.pil_save(str(out_path), format="JPEG", quality=92)
        else:
            pixmap.save(str(out_path))
        outputs.append(str(out_path))

    doc.close()
    return respond(True, output_dir=str(output_dir), outputs=outputs, pages=len(outputs), dpi=args.dpi, format=args.format)


def cmd_export_text(args: argparse.Namespace) -> int:
    import fitz

    output = Path(args.output)
    ensure_parent(output)

    if args.format == "md":
        try:
            capture = io.StringIO()
            with silence_process_output(), contextlib.redirect_stdout(capture), contextlib.redirect_stderr(capture):
                import pymupdf4llm

                text = pymupdf4llm.to_markdown(args.input)
        except Exception:
            doc = fitz.open(args.input)
            chunks = [f"# Page {i + 1}\n\n{page.get_text('text')}" for i, page in enumerate(doc)]
            text = "\n\n".join(chunks)
            doc.close()
    elif args.format == "html":
        doc = fitz.open(args.input)
        chunks = ["<!doctype html><html><head><meta charset=\"utf-8\"><title>PDF Export</title></head><body>"]
        for index, page in enumerate(doc, start=1):
            chunks.append(f"<section data-page=\"{index}\">")
            chunks.append(page.get_text("html"))
            chunks.append("</section>")
        chunks.append("</body></html>")
        text = "\n".join(chunks)
        doc.close()
    else:
        doc = fitz.open(args.input)
        chunks = [f"--- Page {i + 1} ---\n{page.get_text('text')}" for i, page in enumerate(doc)]
        text = "\n\n".join(chunks)
        doc.close()

    output.write_text(text, encoding="utf-8")
    return respond(True, output=str(output), bytes=output.stat().st_size, format=args.format)


def cmd_export_docx(args: argparse.Namespace) -> int:
    from pdf2docx import Converter

    output = Path(args.output)
    ensure_parent(output)
    converter = Converter(args.input)
    try:
        converter.convert(str(output), start=0, end=None)
    finally:
        converter.close()
    return respond(True, output=str(output), bytes=output.stat().st_size)


def cmd_export_xlsx(args: argparse.Namespace) -> int:
    import fitz
    from openpyxl import Workbook

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "PDF Text"
    sheet.append(["Page", "Line", "Text"])

    for page_index, page in enumerate(doc, start=1):
        lines = [line.strip() for line in page.get_text("text").splitlines() if line.strip()]
        for line_index, line in enumerate(lines, start=1):
            sheet.append([page_index, line_index, line])

    doc.close()
    workbook.save(output)
    return respond(True, output=str(output), bytes=output.stat().st_size)


def cmd_export_pptx(args: argparse.Namespace) -> int:
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    presentation = Presentation()
    blank = presentation.slide_layouts[6]

    if doc.page_count > 0:
        first = doc[0].rect
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(max(5.625, 10 * first.height / max(first.width, 1)))

    with tempfile.TemporaryDirectory(prefix="sampdf-pptx-") as tmp_dir:
        for index, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(args.dpi / 72.0, args.dpi / 72.0), alpha=False)
            image_path = Path(tmp_dir) / f"page-{index:03d}.png"
            pix.save(str(image_path))
            slide = presentation.slides.add_slide(blank)
            slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)

    doc.close()
    presentation.save(output)
    return respond(True, output=str(output), bytes=output.stat().st_size)


def cmd_images_to_pdf(args: argparse.Namespace) -> int:
    from PIL import Image, ImageOps

    output = Path(args.output)
    ensure_parent(output)
    images = []
    for image_path in args.input:
        image = Image.open(image_path)
        image = ImageOps.exif_transpose(image)
        if image.mode not in ("RGB", "L"):
            image = image.convert("RGB")
        elif image.mode == "L":
            image = image.convert("RGB")
        images.append(image)

    if not images:
        return respond(False, message="no images supplied")

    first, rest = images[0], images[1:]
    first.save(output, "PDF", resolution=args.dpi, save_all=True, append_images=rest)
    for image in images:
        image.close()
    return respond(True, output=str(output), pages=1 + len(rest), bytes=output.stat().st_size)


def cmd_enhance_scan(args: argparse.Namespace) -> int:
    import fitz
    from PIL import Image, ImageEnhance, ImageFilter

    output = Path(args.output)
    ensure_parent(output)
    doc = fitz.open(args.input)
    images = []
    matrix = fitz.Matrix(args.dpi / 72.0, args.dpi / 72.0)

    with tempfile.TemporaryDirectory(prefix="sampdf-enhance-") as tmp_dir:
        for index, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            if args.grayscale:
                image = image.convert("L").convert("RGB")
            image = ImageEnhance.Contrast(image).enhance(args.contrast)
            image = ImageEnhance.Sharpness(image).enhance(args.sharpness)
            if args.denoise:
                image = image.filter(ImageFilter.MedianFilter(size=3))
            image_path = Path(tmp_dir) / f"page-{index:03d}.jpg"
            image.save(image_path, "JPEG", quality=92)
            saved = Image.open(image_path).convert("RGB")
            images.append(saved)

        if not images:
            return respond(False, message="input has no pages")
        images[0].save(output, "PDF", resolution=args.dpi, save_all=True, append_images=images[1:])

    for image in images:
        image.close()
    doc.close()
    return respond(True, output=str(output), pages=len(images), bytes=output.stat().st_size)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Sam PDF Studio private PDF engine")
    subparsers = parser.add_subparsers(dest="command", required=True)

    health = subparsers.add_parser("health")
    health.set_defaults(func=cmd_health)

    metadata = subparsers.add_parser("metadata")
    metadata.add_argument("--input", required=True)
    metadata.set_defaults(func=cmd_metadata)

    merge = subparsers.add_parser("merge")
    merge.add_argument("--input", action="append", required=True)
    merge.add_argument("--output", required=True)
    merge.set_defaults(func=cmd_merge)

    merge_pages = subparsers.add_parser("merge-pages")
    merge_pages.add_argument("--page-item", action="append", required=True)
    merge_pages.add_argument("--output", required=True)
    merge_pages.set_defaults(func=cmd_merge_pages)

    split = subparsers.add_parser("split")
    split.add_argument("--input", required=True)
    split.add_argument("--output-dir", required=True)
    split.add_argument("--pages")
    split.set_defaults(func=cmd_split)

    extract = subparsers.add_parser("extract-pages")
    extract.add_argument("--input", required=True)
    extract.add_argument("--output", required=True)
    extract.add_argument("--pages", required=True)
    extract.set_defaults(func=cmd_extract_pages)

    delete = subparsers.add_parser("delete-pages")
    delete.add_argument("--input", required=True)
    delete.add_argument("--output", required=True)
    delete.add_argument("--pages", required=True)
    delete.set_defaults(func=cmd_delete_pages)

    rotate = subparsers.add_parser("rotate-pages")
    rotate.add_argument("--input", required=True)
    rotate.add_argument("--output", required=True)
    rotate.add_argument("--pages")
    rotate.add_argument("--degrees", type=int, choices=[90, 180, 270, -90, -180, -270], default=90)
    rotate.set_defaults(func=cmd_rotate_pages)

    crop = subparsers.add_parser("crop-pages")
    crop.add_argument("--input", required=True)
    crop.add_argument("--output", required=True)
    crop.add_argument("--pages")
    crop.add_argument("--left", type=float, default=0)
    crop.add_argument("--top", type=float, default=0)
    crop.add_argument("--right", type=float, default=0)
    crop.add_argument("--bottom", type=float, default=0)
    crop.set_defaults(func=cmd_crop_pages)

    replace = subparsers.add_parser("replace-text")
    replace.add_argument("--input", required=True)
    replace.add_argument("--output", required=True)
    replace.add_argument("--find", required=True)
    replace.add_argument("--replace", required=True)
    replace.add_argument("--font-size", type=float, default=11)
    replace.add_argument("--expand", type=float, default=160)
    replace.add_argument("--auto-size", action="store_true")
    replace.add_argument("--match-style", action=argparse.BooleanOptionalAction, default=True)
    replace.add_argument("--page", type=int)
    replace.add_argument("--rect")
    replace.set_defaults(func=cmd_replace_text)

    page_blocks = subparsers.add_parser("page-blocks")
    page_blocks.add_argument("--input", required=True)
    page_blocks.add_argument("--page", type=int, required=True)
    page_blocks.set_defaults(func=cmd_page_blocks)

    block_at = subparsers.add_parser("block-at")
    block_at.add_argument("--input", required=True)
    block_at.add_argument("--page", type=int, required=True)
    block_at.add_argument("--x", type=float, required=True)
    block_at.add_argument("--y", type=float, required=True)
    block_at.set_defaults(func=cmd_block_at)

    replace_block = subparsers.add_parser("replace-block")
    replace_block.add_argument("--input", required=True)
    replace_block.add_argument("--output", required=True)
    replace_block.add_argument("--page", type=int, required=True)
    replace_block.add_argument("--rect", required=True)
    replace_block.add_argument("--text", required=True)
    replace_block.add_argument("--font-size", type=float, default=0)
    replace_block.add_argument("--font", default=None)
    replace_block.add_argument("--color")
    replace_block.add_argument("--bold", action="store_true")
    replace_block.add_argument("--italic", action="store_true")
    replace_block.add_argument("--underline", action="store_true")
    replace_block.add_argument("--align", choices=["left", "center", "right"], default="left")
    replace_block.add_argument("--background", default=None)
    replace_block.add_argument("--expand", type=float, default=40)
    replace_block.add_argument("--line-height", type=float, default=0)
    replace_block.add_argument("--track-original")
    replace_block.add_argument("--original-text")
    replace_block.add_argument("--color-runs")
    replace_block.set_defaults(func=cmd_replace_block)

    block_background = subparsers.add_parser("block-background")
    block_background.add_argument("--input", required=True)
    block_background.add_argument("--output", required=True)
    block_background.add_argument("--page", type=int, required=True)
    block_background.add_argument("--rect", required=True)
    block_background.add_argument("--color", required=True)
    block_background.add_argument("--padding", type=float, default=3)
    block_background.set_defaults(func=cmd_block_background)

    move_block = subparsers.add_parser("move-block")
    move_block.add_argument("--input", required=True)
    move_block.add_argument("--output", required=True)
    move_block.add_argument("--page", type=int, required=True)
    move_block.add_argument("--rect", required=True)
    move_block.add_argument("--text", required=True)
    move_block.add_argument("--dest-x", type=float, required=True)
    move_block.add_argument("--dest-y", type=float, required=True)
    move_block.add_argument("--font-size", type=float, default=0)
    move_block.add_argument("--font", default=None)
    move_block.add_argument("--color")
    move_block.add_argument("--bold", action="store_true")
    move_block.add_argument("--italic", action="store_true")
    move_block.add_argument("--underline", action="store_true")
    move_block.add_argument("--align", choices=["left", "center", "right"], default="left")
    move_block.add_argument("--expand", type=float, default=40)
    move_block.add_argument("--line-height", type=float, default=0)
    move_block.add_argument("--original-text")
    move_block.set_defaults(func=cmd_move_block)

    add_symbol = subparsers.add_parser("add-symbol")
    add_symbol.add_argument("--input", required=True)
    add_symbol.add_argument("--output", required=True)
    add_symbol.add_argument("--page", type=int, required=True)
    add_symbol.add_argument("--kind", required=True)
    add_symbol.add_argument("--x", type=float, required=True)
    add_symbol.add_argument("--y", type=float, required=True)
    add_symbol.add_argument("--size", type=float, default=16)
    add_symbol.add_argument("--color")
    add_symbol.set_defaults(func=cmd_add_symbol)

    compress = subparsers.add_parser("compress")
    compress.add_argument("--input", required=True)
    compress.add_argument("--output", required=True)
    compress.add_argument("--quality", default="medium")
    compress.set_defaults(func=cmd_compress)

    page_numbers = subparsers.add_parser("add-page-numbers")
    page_numbers.add_argument("--input", required=True)
    page_numbers.add_argument("--output", required=True)
    page_numbers.add_argument("--position", default="bottom-center")
    page_numbers.add_argument("--number-format", default="n")
    page_numbers.add_argument("--start", type=int, default=1)
    page_numbers.add_argument("--font-size", type=float, default=11)
    page_numbers.add_argument("--margin", type=float, default=28)
    page_numbers.set_defaults(func=cmd_add_page_numbers)

    resize_pages = subparsers.add_parser("resize-pages")
    resize_pages.add_argument("--input", required=True)
    resize_pages.add_argument("--output", required=True)
    resize_pages.add_argument("--width", type=float, required=True)
    resize_pages.add_argument("--height", type=float, required=True)
    resize_pages.set_defaults(func=cmd_resize_pages)

    set_password = subparsers.add_parser("set-password")
    set_password.add_argument("--input", required=True)
    set_password.add_argument("--output", required=True)
    set_password.add_argument("--password", required=True)
    set_password.set_defaults(func=cmd_set_password)

    redline = subparsers.add_parser("redline")
    redline.add_argument("--input", required=True)
    redline.add_argument("--output", required=True)
    redline.add_argument("--page", type=int, required=True)
    redline.add_argument("--kind", required=True)
    redline.add_argument("--rects")
    redline.add_argument("--x", type=float, default=0)
    redline.add_argument("--y", type=float, default=0)
    redline.add_argument("--note")
    redline.set_defaults(func=cmd_redline)

    redact = subparsers.add_parser("redact-text")
    redact.add_argument("--input", required=True)
    redact.add_argument("--output", required=True)
    redact.add_argument("--find", required=True)
    redact.add_argument("--label", default="")
    redact.add_argument("--pages")
    redact.add_argument("--page", type=int)
    redact.add_argument("--rect")
    redact.set_defaults(func=cmd_redact_text)

    add_text = subparsers.add_parser("add-text")
    add_text.add_argument("--input", required=True)
    add_text.add_argument("--output", required=True)
    add_text.add_argument("--page", type=int, required=True)
    add_text.add_argument("--x", type=float, required=True)
    add_text.add_argument("--y", type=float, required=True)
    add_text.add_argument("--text", required=True)
    add_text.add_argument("--font-size", type=float, default=12)
    add_text.add_argument("--font", default="Helvetica")
    add_text.add_argument("--color", default="#000000")
    add_text.add_argument("--bold", action="store_true")
    add_text.add_argument("--italic", action="store_true")
    add_text.add_argument("--underline", action="store_true")
    add_text.add_argument("--align", choices=["left", "center", "right"], default="left")
    add_text.set_defaults(func=cmd_add_text)

    annotate = subparsers.add_parser("annotate-text")
    annotate.add_argument("--input", required=True)
    annotate.add_argument("--output", required=True)
    annotate.add_argument("--kind", choices=["highlight", "underline", "strikeout"], required=True)
    annotate.add_argument("--find", required=True)
    annotate.add_argument("--note")
    annotate.add_argument("--pages")
    annotate.set_defaults(func=cmd_annotate_text)

    note = subparsers.add_parser("add-note")
    note.add_argument("--input", required=True)
    note.add_argument("--output", required=True)
    note.add_argument("--page", type=int, required=True)
    note.add_argument("--x", type=float, required=True)
    note.add_argument("--y", type=float, required=True)
    note.add_argument("--text", required=True)
    note.set_defaults(func=cmd_add_note)

    signature = subparsers.add_parser("add-signature")
    signature.add_argument("--input", required=True)
    signature.add_argument("--output", required=True)
    signature.add_argument("--page", type=int, required=True)
    signature.add_argument("--x", type=float, required=True)
    signature.add_argument("--y", type=float, required=True)
    signature.add_argument("--width", type=float, default=180)
    signature.add_argument("--height", type=float, default=42)
    signature.add_argument("--text", required=True)
    signature.add_argument("--font-size", type=float, default=24)
    signature.set_defaults(func=cmd_add_signature)

    add_image = subparsers.add_parser("add-image")
    add_image.add_argument("--input", required=True)
    add_image.add_argument("--output", required=True)
    add_image.add_argument("--image", required=True)
    add_image.add_argument("--page", type=int, required=True)
    add_image.add_argument("--x", type=float, required=True)
    add_image.add_argument("--y", type=float, required=True)
    add_image.add_argument("--width", type=float, required=True)
    add_image.add_argument("--height", type=float, required=True)
    add_image.set_defaults(func=cmd_add_image)

    paste_region = subparsers.add_parser("paste-region")
    paste_region.add_argument("--input", required=True)
    paste_region.add_argument("--source", required=True)
    paste_region.add_argument("--output", required=True)
    paste_region.add_argument("--source-page", type=int, required=True)
    paste_region.add_argument("--source-rect", required=True)
    paste_region.add_argument("--destination-page", type=int, required=True)
    paste_region.add_argument("--destination-x", type=float, required=True)
    paste_region.add_argument("--destination-y", type=float, required=True)
    paste_region.add_argument("--dpi", type=int, default=220)
    paste_region.add_argument("--erase-source", action="store_true")
    paste_region.set_defaults(func=cmd_paste_region)

    link = subparsers.add_parser("add-link")
    link.add_argument("--input", required=True)
    link.add_argument("--output", required=True)
    link.add_argument("--page", type=int, required=True)
    link.add_argument("--x", type=float, required=True)
    link.add_argument("--y", type=float, required=True)
    link.add_argument("--width", type=float, required=True)
    link.add_argument("--height", type=float, required=True)
    link.add_argument("--url", required=True)
    link.set_defaults(func=cmd_add_link)

    ocr = subparsers.add_parser("ocr")
    ocr.add_argument("--input", required=True)
    ocr.add_argument("--output", required=True)
    ocr.add_argument("--language", default="eng")
    ocr.add_argument("--force", action="store_true")
    ocr.add_argument("--page", type=int)
    ocr.set_defaults(func=cmd_ocr)

    enhance = subparsers.add_parser("enhance-scan")
    enhance.add_argument("--input", required=True)
    enhance.add_argument("--output", required=True)
    enhance.add_argument("--dpi", type=int, default=220)
    enhance.add_argument("--contrast", type=float, default=1.35)
    enhance.add_argument("--sharpness", type=float, default=1.15)
    enhance.add_argument("--grayscale", action="store_true")
    enhance.add_argument("--denoise", action="store_true")
    enhance.set_defaults(func=cmd_enhance_scan)

    images = subparsers.add_parser("export-images")
    images.add_argument("--input", required=True)
    images.add_argument("--output-dir", required=True)
    images.add_argument("--format", choices=["png", "jpg"], default="png")
    images.add_argument("--dpi", type=int, default=200)
    images.set_defaults(func=cmd_export_images)

    text = subparsers.add_parser("export-text")
    text.add_argument("--input", required=True)
    text.add_argument("--output", required=True)
    text.add_argument("--format", choices=["txt", "md", "html"], default="txt")
    text.set_defaults(func=cmd_export_text)

    docx = subparsers.add_parser("export-docx")
    docx.add_argument("--input", required=True)
    docx.add_argument("--output", required=True)
    docx.set_defaults(func=cmd_export_docx)

    xlsx = subparsers.add_parser("export-xlsx")
    xlsx.add_argument("--input", required=True)
    xlsx.add_argument("--output", required=True)
    xlsx.set_defaults(func=cmd_export_xlsx)

    pptx = subparsers.add_parser("export-pptx")
    pptx.add_argument("--input", required=True)
    pptx.add_argument("--output", required=True)
    pptx.add_argument("--dpi", type=int, default=160)
    pptx.set_defaults(func=cmd_export_pptx)

    images_pdf = subparsers.add_parser("images-to-pdf")
    images_pdf.add_argument("--input", action="append", required=True)
    images_pdf.add_argument("--output", required=True)
    images_pdf.add_argument("--dpi", type=int, default=200)
    images_pdf.set_defaults(func=cmd_images_to_pdf)

    return parser


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()
    try:
        return args.func(args)
    except Exception as exc:
        return respond(False, message=str(exc), traceback=traceback.format_exc(limit=8))


if __name__ == "__main__":
    raise SystemExit(main())
